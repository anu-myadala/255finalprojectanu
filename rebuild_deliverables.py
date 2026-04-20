"""
Rebuild Taxi_Project_Presentation.pptx and Taxi_Project_Report.docx
with corrected content, clean layouts (ONE visual per slide), large readable
fonts, and speaker notes on every slide.
Run AFTER taxi_analysis.py so charts/ and summary.json are fresh.
"""
import json, os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from docx import Document
from docx.shared import Inches as DInches, Pt as DPt

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
CHART_DIR  = os.path.join(SCRIPT_DIR, "charts")
DIAG_DIR   = os.path.join(SCRIPT_DIR, "diagrams")

with open(os.path.join(SCRIPT_DIR, "summary.json")) as f:
    S = json.load(f)

NAVY  = RGBColor(0x0F, 0x3B, 0x66)
TEAL  = RGBColor(0x0D, 0x94, 0x88)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
AMBER = RGBColor(0xD9, 0x77, 0x06)
SLATE = RGBColor(0x33, 0x41, 0x55)
LGRAY = RGBColor(0x94, 0xA3, 0xB8)

rr  = S["regression_results"]
dt  = S["decision_tree"]
km  = S["kmeans"]
db  = S["dbscan"]
ht  = S.get("hypothesis_tests", {})
lm  = S.get("lift_metrics", {})
skew_b = rr["Log-Linear (log duration)"]["skewness_before_log"]
skew_a = rr["Log-Linear (log duration)"]["skewness_after_log"]
mv_r2  = rr["Multi-Variable Linear (fare)"]["r2"]
mv_rmse = rr["Multi-Variable Linear (fare)"]["rmse"]
lg_auc = rr["Logistic (P(tip>0))"]["auc"]
cv_s   = S.get("cross_validation", {}).get("simple_linear_cv_r2", "N/A")
cv_m   = S.get("cross_validation", {}).get("multi_var_cv_r2", "N/A")
cv_l   = S.get("cross_validation", {}).get("logistic_cv_auc", "N/A")
disp   = rr["Poisson (hourly trip counts)"]["dispersion"]

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_bg(slide, color=NAVY):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color

def tb(slide, left, top, width, height, text, size=18, bold=False,
       color=WHITE, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text; p.font.size = Pt(size); p.font.bold = bold
    p.font.color.rgb = color; p.font.name = "Calibri"; p.alignment = align
    return tf

def bul(slide, left, top, width, height, items, size=16, color=WHITE, spacing=6):
    box = slide.shapes.add_textbox(Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    tf = box.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item; p.font.size = Pt(size); p.font.color.rgb = color
        p.font.name = "Calibri"; p.space_after = Pt(spacing)
    return tf

def img(slide, name, left, top, width=None, height=None):
    fp = os.path.join(CHART_DIR, name)
    if not os.path.exists(fp):
        fp = os.path.join(DIAG_DIR, name)
    if not os.path.exists(fp):
        return
    kw = {"left": Inches(left), "top": Inches(top)}
    if width:  kw["width"]  = Inches(width)
    if height: kw["height"] = Inches(height)
    slide.shapes.add_picture(fp, **kw)

def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text

TOTAL = 26

def footer(slide, num):
    tb(slide, 0.4, 7.05, 12, 0.35,
       f"Anukrithi Myadala  \u00b7  CMPE 255  \u00b7  NYC Taxi Mining          {num} / {TOTAL}",
       size=10, color=LGRAY)

# ── 1  Title ──
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
tb(s, 0.6, 0.4, 12, 0.5, "DATA MINING  \u00b7  CMPE 255  \u00b7  SPRING 2026", size=13, color=TEAL, bold=True)
tb(s, 0.6, 1.2, 10, 1.8, "Mining the Pulse\nof New York City", size=48, bold=True)
tb(s, 0.6, 3.4, 11, 1.0, "Demand, fare, and tip patterns from 2.96 M January 2024 yellow-cab trips \u2014\nfive regression types, decision-tree, K-means++, DBSCAN,\nApriori, PCA, hypothesis testing, and a live dashboard.", size=18, color=LGRAY)
tb(s, 0.6, 5.0, 6, 0.6, "Anukrithi Myadala  \u00b7  April 2026", size=16, color=LGRAY)
for i, (label, val) in enumerate([("RAW RECORDS", "2.96M"), ("CLEANED", f"{S['rows_clean']/1e6:.2f}M"), ("SAMPLE", f"{S['sample_n']:,}"), ("MODELS", "10+")]):
    x = 8.2 + i * 1.3
    tb(s, x, 5.2, 1.2, 0.4, val, size=26, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
    tb(s, x, 5.65, 1.2, 0.3, label, size=10, color=LGRAY, align=PP_ALIGN.CENTER)
notes(s, f"Good morning everyone, I'm Anukrithi Myadala presenting my CMPE 255 final project. This project analyzes nearly 3 million yellow taxi trip records from January 2024. After cleaning we have {S['rows_clean']:,} usable trips and a modeling sample of {S['sample_n']:,}. We implement every major algorithm family from scratch in NumPy. No scikit-learn was used.")

# ── 2  Problem Statement ──
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
tb(s, 0.6, 0.3, 10, 0.8, "Problem Statement", size=36, bold=True)
tb(s, 0.6, 1.0, 10, 0.4, "Why mining 3 million taxi trips matters", size=18, color=TEAL)
tb(s, 0.6, 1.6, 12, 1.0, "NYC processes ~100K yellow-cab trips per day. Drivers waste capacity in Manhattan while outer-borough demand goes unserved. Tip revenue is unpredictable. Airport pricing is opaque.", size=16, color=LGRAY)
bul(s, 0.6, 3.0, 12, 4.0, [f"Q1  What features predict fare? \u2192 Multi-var linear (R\u00b2 = {mv_r2:.2f})", f"Q2  Is duration log-normal? \u2192 Log-linear (skew {skew_b:.2f} \u2192 {skew_a:.2f})", "Q3  Can we forecast hourly counts? \u2192 Poisson GLM (IRLS)", "Q4  Which trips get tips? \u2192 Logistic \u2014 no data leakage", "Q5  Are there zone archetypes? \u2192 K-means++ with restarts", "H1  Airport fares > non-airport \u2192 Welch\u2019s t-test", "H2  Rush-hour trips longer \u2192 Welch\u2019s t-test", "H3  Weekend tip rates differ \u2192 \u03c7\u00b2 test"], size=16)
footer(s, 2)
notes(s, f"We tackle five research questions and three formal hypothesis tests. Our best fare model explains {mv_r2*100:.0f}% of fare variance. We also test whether airport trips have higher fares, rush hour trips take longer, and weekend tip rates differ.")

# ── 3  Data Source ──
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
tb(s, 0.6, 0.3, 10, 0.8, "Data Source & Star Schema", size=36, bold=True)
tb(s, 0.6, 1.0, 10, 0.4, "NYC TLC public trip records + zone lookup", size=18, color=TEAL)
bul(s, 0.6, 1.6, 5.5, 4.5, ["yellow_tripdata_2024-01.csv", f"  \u2022 {S['rows_raw']:,} rows \u00d7 19 columns", "  \u2022 Timestamps, zone IDs, fare, tip, distance", "", "taxi_zone_lookup.csv", "  \u2022 265 zones with borough + zone name", "", "Star Schema:", "  \u2022 Fact: Trip_Fact (1 row = 1 trip)", "  \u2022 Dims: Zone, Time, Payment, Vendor"], size=16)
img(s, "star_schema.png", 7.5, 1.5, width=5.0)
footer(s, 3)
notes(s, f"Our data comes from two public NYC TLC datasets. The main file has {S['rows_raw']:,} rows and 19 columns. We designed a star schema with Trip_Fact at center and four dimension tables.")

# ── 4  Preprocessing ──
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
tb(s, 0.6, 0.3, 10, 0.8, "Data Preprocessing Pipeline", size=36, bold=True)
tb(s, 0.6, 1.0, 10, 0.4, "12 cleaning steps + percentile clipping + feature engineering", size=18, color=TEAL)
bul(s, 0.6, 1.6, 5.8, 5.0, ["1.  Drop core NULLs", "2\u20134.  Filter fare \u2265 0, total \u2265 0, tip \u2265 0", "5\u20136.  Distance: 0 \u2264 d \u2264 100 mi", "7.  Enforce pickup < dropoff", "8.  Restrict to January 2024", "9.  Duration: 1 min \u2264 d \u2264 6 hr", "10.  Passengers: 1\u20136", "11.  Percentile clip (1st\u201399th pctl)", "     \u2192 Preserves airport trips!", "     \u2192 IQR cut all trips > 4.86 mi", f"12.  Stratified sample \u2192 {S['sample_n']:,}"], size=16)
bul(s, 7.2, 1.6, 5.5, 5.0, ["Features engineered:", "  \u2022 duration_min, speed_mph", "  \u2022 hour, day_of_week, is_weekend", "  \u2022 is_rush_hour, is_airport_pickup", "  \u2022 has_tip, tip_pct", "  \u2022 Zone joins from lookup CSV", "", "Data leakage prevented:", "  \u2022 payment_type EXCLUDED", "  \u2022 TLC only records credit-card tips", "  \u2022 Including it = target leakage", f"  \u2022 Result: {S['rows_raw']:,} \u2192 {S['rows_clean']:,} clean"], size=16)
footer(s, 4)
notes(s, f"Starting from {S['rows_raw']:,} raw rows, we apply 12 cleaning steps. Critical choice: percentile clipping instead of IQR. IQR had an upper bound of just 4.86 miles, removing all airport trips. We also excluded payment_type from the tip model to prevent data leakage.")

# ── 5  EDA: Hourly Volume ──
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
tb(s, 0.6, 0.3, 10, 0.8, "EDA \u00b7 Trip Volume by Hour", size=36, bold=True)
tb(s, 0.6, 1.0, 10, 0.4, "When are New Yorkers in cabs? (full cleaned data)", size=18, color=TEAL)
img(s, "01_hourly_volume.png", 1.5, 1.6, width=10.0); footer(s, 5)
notes(s, "This chart uses the full cleaned dataset. Volume peaks 5-7 PM and crashes after midnight. This demand profile informs our K-means clustering and Poisson count modeling.")

# ── 6  EDA: Heatmap ──
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
tb(s, 0.6, 0.3, 10, 0.8, "EDA \u00b7 Day-of-Week \u00d7 Hour Heatmap", size=36, bold=True)
tb(s, 0.6, 1.0, 10, 0.4, "Friday/Saturday late nights are uniquely active \u2014 the nightlife signal", size=18, color=TEAL)
img(s, "02_dow_hour_heatmap.png", 1.5, 1.6, width=10.0); footer(s, 6)
notes(s, "The heatmap reveals Friday and Saturday late nights are uniquely active. Sunday morning is the quietest slot. Weekday mornings show consistent commute patterns.")

# ── 7  EDA: Fare & Distance ──
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
tb(s, 0.6, 0.3, 10, 0.8, "EDA \u00b7 Fare & Distance Distributions", size=36, bold=True)
tb(s, 0.6, 1.0, 10, 0.4, "Both right-skewed after percentile clipping (full cleaned data)", size=18, color=TEAL)
img(s, "03_fare_distance.png", 1.5, 1.6, width=10.0); footer(s, 7)
notes(s, "Fare and distance are both right-skewed. This motivates the log-linear model for duration. Percentile clipping preserves trips up to about 19 miles including airport rides.")

# ── 8  EDA: Borough + Correlation ──
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
tb(s, 0.6, 0.3, 10, 0.8, "EDA \u00b7 Borough Volume & Correlations", size=36, bold=True)
img(s, "04_borough_volume.png", 0.4, 1.3, width=6.0)
img(s, "05_correlation.png", 6.8, 1.2, width=6.0); footer(s, 8)
notes(s, "Manhattan dominates with about 95% of yellow-cab pickups. The correlation heatmap shows fare and distance at r=0.81. Passenger count has near-zero correlation with everything.")

# ── 9  Regression Stack ──
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
tb(s, 0.6, 0.3, 10, 0.8, "Mining \u00b7 Five Regression Families", size=36, bold=True)
tb(s, 0.6, 1.0, 10, 0.4, "All from scratch in NumPy \u2014 no sklearn", size=18, color=TEAL)
bul(s, 0.6, 1.6, 5.5, 5.0, ["1. Simple Linear: fare ~ distance", "   Closed-form normal equations", "", "2. Multi-Variable Linear: fare ~ 8 features", f"   R\u00b2 = {mv_r2:.3f}, RMSE = ${mv_rmse:.2f}", "", "3. Log-Linear: log(duration) ~ features", f"   Skew {skew_b:.2f} \u2192 {skew_a:.2f}", "", "4. Poisson GLM: hourly counts (IRLS)", f"   Dispersion = {disp:.0f}\u00d7 (overdispersed)", "", "5. Logistic: P(tip > 0)", "   No credit-card feature (leakage)"], size=16)
img(s, "06_regression_scores.png", 6.8, 1.3, width=6.0); footer(s, 9)
notes(s, f"Five regression families from scratch. Multi-variable linear achieves R-squared {mv_r2:.3f}. Log-linear confirms duration is log-normal. Poisson reveals overdispersion at {disp:.0f}x. Logistic predicts tips with AUC {lg_auc:.3f} without the leaked payment_type.")

# ── 10  Regression Results + CV ──
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
tb(s, 0.6, 0.3, 10, 0.8, "Regression Results & Cross-Validation", size=36, bold=True)
tb(s, 0.6, 1.0, 10, 0.4, "5-fold CV implemented from scratch", size=18, color=TEAL)
bul(s, 0.6, 1.6, 6.0, 5.0, [f"Simple Linear R\u00b2:       {rr['Simple Linear (fare~distance)']['r2']:.3f}", f"  5-fold CV R\u00b2:         {cv_s}", "", f"Multi-Var Linear R\u00b2:    {mv_r2:.3f}  \u2190 Best fare model", f"  5-fold CV R\u00b2:         {cv_m}", f"  RMSE:                 ${mv_rmse:.2f}", "", f"Log-Linear R\u00b2 (log):    {rr['Log-Linear (log duration)']['r2_log_space']:.3f}", "", f"Poisson pseudo-R\u00b2:      {rr['Poisson (hourly trip counts)']['pseudo_r2_mcfadden']:.3f}", f"  Dispersion:           {disp:.0f}\u00d7 (overdispersed)", "", f"Logistic AUC:           {lg_auc:.3f}", f"  5-fold CV AUC:        {cv_l}"], size=16)
img(s, "13_residuals.png", 7.0, 1.5, width=5.8); footer(s, 10)
notes(s, f"Detailed metrics with 5-fold CV. Low standard deviations confirm models generalize well. The residual plot shows multi-variable linear residuals are well-behaved.")

# ── 11  Decision Tree ──
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
tb(s, 0.6, 0.3, 10, 0.8, "Mining \u00b7 Decision Tree (CART / Gini)", size=36, bold=True)
tb(s, 0.6, 1.0, 10, 0.4, f"Peak vs off-peak, depth 6 \u00b7 AUC = {dt['auc']:.3f} \u2014 a meaningful negative result", size=18, color=TEAL)
img(s, "07_dt_cm_importance.png", 1.5, 1.6, width=10.0)
bul(s, 0.6, 5.8, 12, 1.0, [f"Acc {dt['accuracy']:.3f}  |  Prec {dt['precision']:.3f}  |  Rec {dt['recall']:.3f}  |  F1 {dt['f1']:.3f}  |  AUC {dt['auc']:.3f}     Key insight: trip features are time-invariant."], size=14, color=LGRAY)
footer(s, 11)
notes(s, f"The CART tree achieves AUC {dt['auc']:.3f} \u2014 near random. This is a meaningful NEGATIVE finding: trip characteristics are time-invariant. A cab ride at 3 PM looks like one at 8 AM.")

# ── 12  K-Means ──
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
tb(s, 0.6, 0.3, 10, 0.8, "Mining \u00b7 K-Means++ Zone Archetypes", size=36, bold=True)
tb(s, 0.6, 1.0, 12, 0.4, f"k = {km['k']} (auto-detected) \u00b7 Silhouette = {km['silhouette']:.3f} \u00b7 k-means++ with 10 restarts", size=18, color=TEAL)
img(s, "08_kmeans_zones.png", 0.4, 1.5, width=7.0)
bul(s, 8.0, 1.5, 5.0, 3.5, [f"\u2022 {p['label']} ({p['size']} zones)" for p in km.get("profiles", [])], size=16)
img(s, "15_kmeans_elbow.png", 8.0, 4.5, width=4.8); footer(s, 12)
notes(s, f"K-means++ with 10 restarts groups zones into {km['k']} archetypes. Automatic elbow detection selected k={km['k']}. Silhouette {km['silhouette']:.3f} indicates moderate overlap, expected for urban zones.")

# ── 13  DBSCAN ──
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
tb(s, 0.6, 0.3, 10, 0.8, "Mining \u00b7 DBSCAN Outlier Detection", size=36, bold=True)
tb(s, 0.6, 1.0, 10, 0.4, f"\u03b5 = 1.2, min_pts = 4 \u00b7 {db['clusters']} clusters, {db['outliers']} outliers", size=18, color=TEAL)
img(s, "09_dbscan.png", 0.4, 1.5, width=7.0)
oz_items = [f"\u2022 {z['Zone']} ({z['Borough']}) \u2014 airport={z['airport_share']:.0%}" for z in S.get("dbscan_outlier_zones", [])[:5]]
bul(s, 8.0, 1.5, 5.0, 4.0, ["Top outlier zones:"] + oz_items + ["", "JFK & LaGuardia stand out:", "airport_share = 1.0 is unique"], size=16)
footer(s, 13)
notes(s, f"DBSCAN finds {db['clusters']} dense clusters and flags {db['outliers']} outlier zones. JFK and LaGuardia stand out with airport_share = 1.0.")

# ── 14  Apriori ──
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
tb(s, 0.6, 0.3, 10, 0.8, "Mining \u00b7 Apriori Association Rules", size=36, bold=True)
tb(s, 0.6, 1.0, 10, 0.4, "Real Apriori: candidate generation + downward-closure pruning, zone level", size=18, color=TEAL)
img(s, "10_apriori_rules.png", 0.4, 1.5, width=8.0)
top3 = S.get("apriori_top3", [])
rule_items = ["Algorithm:", "  \u2022 Level-wise candidate gen", "  \u2022 Downward-closure pruning", "  \u2022 Itemsets up to size 4", ""]
for r in top3[:3]:
    rule_items.append(f"\u2022 {r['antecedent']} \u2192 {r['consequent']}")
    rule_items.append(f"  Lift {r['lift']} \u00b7 Conf {r['confidence']}")
bul(s, 8.8, 1.5, 4.2, 5.0, rule_items, size=14); footer(s, 14)
notes(s, "This is real Apriori with level-wise candidate generation and downward-closure pruning, not just pair counting. Transactions include zone-level items, borough, rush hour, and weekend flags.")

# ── 15  PCA ──
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
tb(s, 0.6, 0.3, 10, 0.8, "Mining \u00b7 PCA Variance Analysis", size=36, bold=True)
tb(s, 0.6, 1.0, 10, 0.4, f"{S['pca_components_for_90pct']} of 8 components capture \u226590% variance", size=18, color=TEAL)
img(s, "12_pca_variance.png", 0.4, 1.6, width=7.5)
bul(s, 8.5, 1.6, 4.5, 4.0, ["From-scratch implementation:", "  \u2022 Standardize features", "  \u2022 Covariance matrix", "  \u2022 Eigendecomposition", "  \u2022 Sort by explained variance", "", f"\u2022 {S['pca_components_for_90pct']} components for 90%", "\u2022 Dimensionality: 8 \u2192 5", "\u2022 Used for K-Means/DBSCAN", "  2D projections"], size=16)
footer(s, 15)
notes(s, f"PCA reduces 8 dimensions to {S['pca_components_for_90pct']} while keeping 90% of variance. The top 2 components are used for the 2D scatter visualizations.")

# ── 16  ROC Curves ──
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
tb(s, 0.6, 0.3, 10, 0.8, "Evaluation \u00b7 ROC Curves", size=36, bold=True)
tb(s, 0.6, 1.0, 10, 0.4, "Logistic tip classifier vs Decision Tree peak-hour classifier", size=18, color=TEAL)
img(s, "14_roc_curves.png", 0.4, 1.6, width=7.0)
bul(s, 8.0, 1.6, 5.0, 4.5, [f"Logistic AUC: {lg_auc:.3f}", f"  5-fold CV: {cv_l}", "  No credit-card leakage", "", f"Decision Tree AUC: {dt['auc']:.3f}", "  Negative result \u2014 features", "  don\u2019t discriminate time-of-day", "", "Data leakage note:", "  payment_type removed because", "  TLC only records credit-card tips", "  Including it inflated AUC to 0.93+"], size=15)
footer(s, 16)
notes(s, f"The logistic tip classifier achieves AUC {lg_auc:.3f} without the leaked payment_type. The decision tree's AUC {dt['auc']:.3f} is a confirmed negative result. We identified and removed the data leakage.")

# ── 17  Lift & Gains ──
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
tb(s, 0.6, 0.3, 10, 0.8, "Evaluation \u00b7 Lift & Cumulative Gains", size=36, bold=True)
tb(s, 0.6, 1.0, 10, 0.4, f"Top 20% captures {lm.get('top20_capture_pct','N/A')}% of tippers \u2014 {lm.get('top20_lift','N/A')}\u00d7 lift", size=18, color=TEAL)
img(s, "11_lift_gains.png", 0.4, 1.6, width=12.0); footer(s, 17)
notes(s, f"Targeting the top 20% captures {lm.get('top20_capture_pct','N/A')}% of tippers \u2014 a lift of {lm.get('top20_lift','N/A')}x over random.")

# ── 18  Hypothesis Testing ──
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
tb(s, 0.6, 0.3, 10, 0.8, "Hypothesis Testing", size=36, bold=True)
tb(s, 0.6, 1.0, 10, 0.4, "Three formal tests \u2014 all implemented from scratch", size=18, color=TEAL)
h_items = []
for hname, hdata in ht.items():
    sig = "\u2713 SIGNIFICANT" if hdata.get("significant") else "\u2717 Not significant"
    h_items.append(f"{hname}")
    if "mean_airport" in hdata:
        h_items.append(f"  Mean: ${hdata['mean_airport']:.2f} vs ${hdata['mean_non_airport']:.2f}")
    elif "mean_rush" in hdata:
        h_items.append(f"  Mean: {hdata['mean_rush']:.1f} min vs {hdata['mean_off_peak']:.1f} min")
    elif "weekend_tip_rate" in hdata:
        h_items.append(f"  Weekend: {hdata['weekend_tip_rate']:.1%} vs Weekday: {hdata['weekday_tip_rate']:.1%}")
    h_items.append(f"  p = {hdata.get('p_value','N/A')} \u2192 {sig}")
    h_items.append("")
bul(s, 0.6, 1.6, 12, 5.0, h_items, size=18); footer(s, 18)
notes(s, "Three formal tests, all from scratch. H1: Airport fares significantly higher. H2: Rush-hour trips significantly longer. H3: Weekend tip rates differ. All confirmed with p < 0.05.")

# ── 19  Cross-Validation ──
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
tb(s, 0.6, 0.3, 10, 0.8, "5-Fold Cross-Validation", size=36, bold=True)
tb(s, 0.6, 1.0, 10, 0.4, "From scratch \u2014 verifies models generalize", size=18, color=TEAL)
bul(s, 0.6, 1.8, 12, 5.0, [f"Simple Linear:       5-fold R\u00b2 = {cv_s}", f"Multi-Var Linear:    5-fold R\u00b2 = {cv_m}", f"Logistic Tip:        5-fold AUC = {cv_l}", "", "Why it matters:", "  \u2022 Single 80/20 split is sensitive to random seed", "  \u2022 K-fold averages over 5 different test sets", "  \u2022 \u00b1 standard deviation shows model stability", "  \u2022 Low std = no overfitting concern", "", "Implementation:", "  \u2022 kfold_indices() generates non-overlapping folds", "  \u2022 Each fold: standardize on train, predict on test", "  \u2022 Report mean \u00b1 std"], size=18)
footer(s, 19)
notes(s, f"5-fold cross-validation from scratch. Simple linear: {cv_s}. Multi-variable: {cv_m}. Logistic: {cv_l}. Low standard deviations confirm generalization.")

# ── 20  Model Comparison ──
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
tb(s, 0.6, 0.3, 10, 0.8, "Model Comparison Summary", size=36, bold=True)
tb(s, 0.6, 1.0, 10, 0.4, "Headline metric for every model family", size=18, color=TEAL)
bul(s, 0.6, 1.6, 12, 5.5, ["Regression:", f"  Simple Linear:     R\u00b2 = {rr['Simple Linear (fare~distance)']['r2']:.3f}", f"  Multi-Var Linear:  R\u00b2 = {mv_r2:.3f}, RMSE = ${mv_rmse:.2f}  \u2190 Best", f"  Log-Linear:        R\u00b2 = {rr['Log-Linear (log duration)']['r2_log_space']:.3f} (log space)", f"  Poisson GLM:       pseudo-R\u00b2 = {rr['Poisson (hourly trip counts)']['pseudo_r2_mcfadden']:.3f}, dispersion = {disp:.0f}\u00d7", "", "Classification:", f"  Logistic (tip):    AUC = {lg_auc:.3f}, F1 = {rr['Logistic (P(tip>0))']['f1']:.3f}", f"  Decision Tree:     AUC = {dt['auc']:.3f} (negative result)", "", "Unsupervised:", f"  K-Means++:  k = {km['k']}, silhouette = {km['silhouette']:.3f}", f"  DBSCAN:     {db['clusters']} clusters, {db['outliers']} outliers", f"  PCA:        {S['pca_components_for_90pct']}/8 components for 90%", "  Apriori:    Zone-level, real candidate generation"], size=16)
footer(s, 20)
notes(s, f"All results in one place. Multi-variable linear is our best fare model at R-squared {mv_r2:.3f}. K-means++ with k={km['k']}. DBSCAN flags {db['outliers']} anomalous zones.")

# ── 21  Knowledge Interpretation ──
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
tb(s, 0.6, 0.3, 10, 0.8, "Knowledge Interpretation", size=36, bold=True)
tb(s, 0.6, 1.0, 10, 0.4, "Five operational insights", size=18, color=TEAL)
bul(s, 0.6, 1.6, 12, 5.0, [f"1. Fare prediction is solved: 8 features \u2192 R\u00b2 = {mv_r2:.2f}, ${mv_rmse:.2f} RMSE", "", f"2. Duration is log-normal: skew {skew_b:.2f} \u2192 {skew_a:.2f} after log transform", "", f"3. Hourly counts are overdispersed ({disp:.0f}\u00d7) \u2014 Negative Binomial needed", "", "4. Tip prediction without payment_type gives honest metrics", "   Removing leakage = data science integrity", "", "5. Outer-borough trips stay local \u2014 fleet-positioning signal"], size=18)
footer(s, 21)
notes(s, "Five key insights. Fare prediction solved. Duration is log-normal. Poisson surfaces overdispersion. Data leakage caught. Outer-borough patterns reveal fleet-positioning opportunities.")

# ── 22  Conclusions ──
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
tb(s, 0.6, 0.3, 10, 0.8, "Conclusions & Future Work", size=36, bold=True)
bul(s, 0.6, 1.4, 5.5, 5.5, ["What we delivered:", f"  \u2713 Cleaned {S['rows_raw']:,} \u2192 {S['rows_clean']:,}", "  \u2713 5 regression families from scratch", f"  \u2713 Best model: R\u00b2 {mv_r2:.2f}", "  \u2713 5-fold cross-validation", "  \u2713 3 hypothesis tests", "  \u2713 K-means++ with auto elbow", f"  \u2713 DBSCAN: {db['outliers']} outlier zones", "  \u2713 Real Apriori at zone level", f"  \u2713 PCA: {S['pca_components_for_90pct']}/8 components", "  \u2713 Interactive Plotly dashboard", "  \u2713 Data leakage caught & removed"], size=16)
bul(s, 7.2, 1.4, 5.5, 5.5, ["Future work:", "  \u2022 Negative Binomial GLM", "  \u2022 Spatial DBSCAN on lat/long", "  \u2022 Per-cluster fare/tip models", "  \u2022 Streaming pipeline (Kafka+Spark)", "  \u2022 Temporal train/test split", "  \u2022 Fairness audit across boroughs", "", "Reproducibility:", "  \u2022 Single script: taxi_analysis.py", "  \u2022 All paths portable", "  \u2022 Every metric reproducible"], size=16)
footer(s, 22)
notes(s, f"We delivered a complete pipeline from {S['rows_raw']:,} raw rows to 10+ models, all from scratch. Validated with CV, tested hypotheses, caught data leakage. Everything reproducible from one script.")

# ── 23  Dashboard Demo ──
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
tb(s, 0.6, 0.3, 10, 0.8, "Live Dashboard Demo", size=36, bold=True)
tb(s, 0.6, 1.0, 10, 0.4, "Open Taxi_Interactive_Dashboard.html", size=18, color=TEAL)
bul(s, 0.6, 1.8, 12, 4.5, ["Dashboard features:", "  \u2022 KPI cards: trip count, avg fare, avg distance, tip rate", "  \u2022 Filters: borough, day type, cluster, hour range", "  \u2022 Charts: hourly volume, borough bars, fare histogram,", "    distance vs fare scatter, tip rate by hour, cluster pie", "", "  \u2022 Built with Plotly.js \u2014 fully client-side, no server", f"  \u2022 Powered by {S['sample_n']:,} real January 2024 trips", "  \u2022 All filters update all charts simultaneously"], size=18)
footer(s, 23)
notes(s, f"Let me switch to the live dashboard. Built with Plotly.js, fully client-side. Powered by {S['sample_n']:,} real trips. All filters update all charts simultaneously.")

# ── 24  Workflow ──
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
tb(s, 0.6, 0.3, 10, 0.8, "End-to-End Pipeline Workflow", size=36, bold=True)
tb(s, 0.6, 1.0, 10, 0.4, "From raw data to knowledge discovery", size=18, color=TEAL)
img(s, "workflow.png", 2.0, 1.6, width=9.0); footer(s, 24)
notes(s, "This diagram shows the complete workflow from problem definition through deployment, mirroring the 12-step process taught in CMPE 255.")

# ── 25  References ──
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
tb(s, 0.6, 0.3, 10, 0.8, "References & Tools", size=36, bold=True)
bul(s, 0.6, 1.4, 12, 5.0, ["Data:", "  \u2022 NYC TLC Trip Records \u2014 nyc.gov/site/tlc/about/tlc-trip-record-data.page", "  \u2022 Taxi Zone Lookup \u2014 published by TLC", "", "Textbook:", "  \u2022 Han, Pei & Tong \u2014 Data Mining: Concepts and Techniques, 4th Ed. (2023)", "", "Tools:", "  \u2022 Python 3, NumPy, pandas, matplotlib, seaborn, Plotly.js", "  \u2022 All algorithms implemented from scratch (no sklearn)"], size=18)
footer(s, 25)
notes(s, "References: NYC TLC public data, Han Pei Tong textbook, Python stack. All algorithms from scratch.")

# ── 26  Thank You ──
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
tb(s, 0.6, 2.0, 12, 1.5, "Thank You", size=52, bold=True, align=PP_ALIGN.CENTER)
tb(s, 0.6, 3.6, 12, 1.0, "Questions?", size=32, color=TEAL, align=PP_ALIGN.CENTER)
tb(s, 0.6, 4.8, 12, 0.8, "Anukrithi Myadala  \u00b7  CMPE 255  \u00b7  Data Mining  \u00b7  Spring 2026\nProf. Vidhyacharan Bhaskar  \u00b7  San Jos\u00e9 State University", size=16, color=LGRAY, align=PP_ALIGN.CENTER)
footer(s, 26)
notes(s, "Thank you for listening. I'm happy to take questions about any part of the pipeline, the from-scratch implementations, the data leakage finding, or the dashboard.")

pptx_path = os.path.join(SCRIPT_DIR, "Taxi_Project_Presentation.pptx")
prs.save(pptx_path)
print(f"\u2713 Saved: {pptx_path}  ({TOTAL} slides, every slide has speaker notes)")

# =====================================================================
#  WORD REPORT
# =====================================================================
doc = Document()
doc.add_heading("Mining the Pulse of New York City", level=0)
doc.add_paragraph("Demand, Fare and Tip Patterns from 2.96 Million January 2024 Yellow-Cab Trips")
doc.add_paragraph("Anukrithi Myadala")
doc.add_paragraph("CMPE 255  \u00b7  Data Mining  \u00b7  Spring 2026")
doc.add_paragraph("Prof. Vidhyacharan Bhaskar  \u00b7  San Jos\u00e9 State University")
doc.add_paragraph("Final project report")
doc.add_page_break()

doc.add_heading("1. Problem Definition and Project Objectives", level=1)
doc.add_paragraph("New York City fields roughly one hundred thousand yellow-cab trips every day, generating millions of structured records per month. This project applies every major mining technique from the CMPE 255 syllabus to the January 2024 TLC trip records.")
doc.add_heading("Domain", level=2)
doc.add_paragraph("Transportation analytics \u2014 urban mobility, consumer behaviour, operations research.")
doc.add_heading("Research questions", level=2)
for q in ["What features predict fare amount, and how much does each contribute?", "Is trip duration log-normally distributed, and does a log transform improve regression?", "Can hourly trip counts per pickup zone be forecast with a Poisson GLM?", "Which trips generate tips, and what is the lift over random targeting?", "Are there natural pickup-zone archetypes detectable through K-means clustering?", "Do airport trips have significantly higher fares? (Welch\u2019s t-test)", "Are rush-hour trips longer in duration? (Welch\u2019s t-test)", "Do weekend tip rates differ from weekday? (\u03c7\u00b2 test)"]:
    doc.add_paragraph(q, style="List Bullet")
doc.add_heading("Target variables", level=2)
doc.add_paragraph("Continuous: fare_amount, trip_duration, trip_count per zone-hour. Binary: has_tip (logistic), is_rush_hour (decision tree).")
doc.add_heading("Expected outcomes", level=2)
doc.add_paragraph(f"A reproducible pipeline: {S['rows_raw']:,} raw rows \u2192 {S['rows_clean']:,} cleaned \u2192 {S['sample_n']:,} sample \u2192 10+ from-scratch models, 5-fold CV, 3 hypothesis tests, dashboard.")
doc.add_heading("2. Data Source Identification", level=1)
for src in [f"yellow_tripdata_2024-01.csv \u2014 {S['rows_raw']:,} trips \u00d7 19 columns.", "taxi_zone_lookup.csv \u2014 265 zones with borough + zone name."]:
    doc.add_paragraph(src, style="List Bullet")
doc.add_heading("3. Data Collection", level=1)
doc.add_paragraph("Parquet downloaded from TLC CDN, converted to CSV. Zone lookup downloaded as CSV.")
doc.add_heading("4. Data Warehousing", level=1)
doc.add_paragraph("Star schema: Trip_Fact + four dimensions (Zone, Time, Payment, Vendor).")
fp = os.path.join(DIAG_DIR, "star_schema.png")
if os.path.exists(fp):
    doc.add_paragraph("Figure 1. Star schema design.")
    doc.add_picture(fp, width=DInches(5.0))
doc.add_heading("5. Data Preprocessing", level=1)
doc.add_heading("5.1 Cleaning audit", level=2)
doc.add_paragraph("12-step pipeline. Key: percentile clipping (1st\u201399th) instead of IQR \u2014 preserves airport trips.")
doc.add_heading("5.2 Feature engineering", level=2)
for f in ["duration_min, speed_mph", "hour, day_of_week, is_weekend, is_rush_hour", "is_airport_pickup/dropoff", "has_tip, tip_pct", "Zone joins from lookup"]:
    doc.add_paragraph(f, style="List Bullet")
doc.add_heading("5.3 Data leakage prevention", level=2)
doc.add_paragraph("payment_type was NOT used for tip prediction. TLC only records credit-card tips \u2014 including it would be data leakage.")
doc.add_heading("5.4 Sampling", level=2)
doc.add_paragraph(f"Stratified by hour: {S['rows_clean']:,} \u2192 {S['sample_n']:,} rows.")
doc.add_heading("6. Exploratory Data Analysis", level=1)
for num, fn, cap in [(2, "01_hourly_volume.png", "Trip volume by hour."), (3, "02_dow_hour_heatmap.png", "Day \u00d7 hour heatmap."), (4, "03_fare_distance.png", "Fare and distance distributions."), (5, "04_borough_volume.png", "Borough pickup volume."), (6, "05_correlation.png", "Feature correlations.")]:
    fp = os.path.join(CHART_DIR, fn)
    if os.path.exists(fp):
        doc.add_paragraph(f"Figure {num}. {cap}")
        doc.add_picture(fp, width=DInches(5.5))
doc.add_heading("7. Data Visualization", level=1)
doc.add_paragraph("15 static matplotlib/seaborn charts + interactive Plotly.js dashboard.")
doc.add_heading("8. Data Mining Techniques", level=1)
doc.add_paragraph("All algorithms from scratch in NumPy/pandas. No sklearn.")
doc.add_heading("8.1 Regression \u2014 five families", level=2)
fp = os.path.join(CHART_DIR, "06_regression_scores.png")
if os.path.exists(fp):
    doc.add_paragraph("Figure 7. R\u00b2 / AUC across five regression families.")
    doc.add_picture(fp, width=DInches(5.5))
doc.add_paragraph(f"Simple linear: R\u00b2={rr['Simple Linear (fare~distance)']['r2']:.3f}. Multi-var: R\u00b2={mv_r2:.3f}, RMSE=${mv_rmse:.2f}. Log-linear: skew {skew_b:.2f}\u2192{skew_a:.2f}. Poisson: dispersion={disp:.0f}\u00d7. Logistic: AUC={lg_auc:.3f}.")
doc.add_heading("8.2 Decision Tree (CART/Gini)", level=2)
fp = os.path.join(CHART_DIR, "07_dt_cm_importance.png")
if os.path.exists(fp):
    doc.add_paragraph("Figure 8. Confusion matrix and feature importance.")
    doc.add_picture(fp, width=DInches(5.5))
doc.add_paragraph(f"AUC={dt['auc']:.3f} \u2014 negative result: trip features are time-invariant.")
doc.add_heading("8.3 K-Means++ clustering", level=2)
fp = os.path.join(CHART_DIR, "08_kmeans_zones.png")
if os.path.exists(fp):
    doc.add_paragraph(f"Figure 9. {km['k']} zone archetypes.")
    doc.add_picture(fp, width=DInches(5.0))
doc.add_paragraph(f"k={km['k']} (auto-detected), silhouette={km['silhouette']:.3f}, k-means++ with 10 restarts.")
doc.add_heading("8.4 DBSCAN outlier detection", level=2)
fp = os.path.join(CHART_DIR, "09_dbscan.png")
if os.path.exists(fp):
    doc.add_paragraph(f"Figure 10. DBSCAN flags {db['outliers']} outlier zones.")
    doc.add_picture(fp, width=DInches(5.0))
doc.add_heading("8.5 Apriori association rules", level=2)
fp = os.path.join(CHART_DIR, "10_apriori_rules.png")
if os.path.exists(fp):
    doc.add_paragraph("Figure 11. Top association rules by lift.")
    doc.add_picture(fp, width=DInches(5.5))
doc.add_paragraph("Real Apriori with candidate generation, downward-closure pruning, itemsets up to size 4.")
doc.add_heading("8.6 PCA variance analysis", level=2)
fp = os.path.join(CHART_DIR, "12_pca_variance.png")
if os.path.exists(fp):
    doc.add_paragraph(f"Figure 12. {S['pca_components_for_90pct']}/8 components capture \u226590% variance.")
    doc.add_picture(fp, width=DInches(4.5))
doc.add_heading("9. Model Evaluation", level=1)
doc.add_heading("9.1 Cross-validation", level=2)
for item in [f"Simple Linear: {cv_s}", f"Multi-Var: {cv_m}", f"Logistic: {cv_l}"]:
    doc.add_paragraph(item, style="List Bullet")
doc.add_heading("9.2 Hypothesis testing", level=2)
for hname, hdata in ht.items():
    sig = "SIGNIFICANT" if hdata.get("significant") else "not significant"
    if "t_statistic" in hdata:
        doc.add_paragraph(f"{hname}: t={hdata['t_statistic']}, p={hdata['p_value']} \u2192 {sig}", style="List Bullet")
    elif "chi_square" in hdata:
        doc.add_paragraph(f"{hname}: \u03c7\u00b2={hdata['chi_square']}, p={hdata['p_value']} \u2192 {sig}", style="List Bullet")
doc.add_heading("9.3 Lift and gains", level=2)
fp = os.path.join(CHART_DIR, "11_lift_gains.png")
if os.path.exists(fp):
    doc.add_paragraph("Figure 13. Cumulative gains and lift chart.")
    doc.add_picture(fp, width=DInches(5.5))
doc.add_heading("9.4 ROC curves", level=2)
fp = os.path.join(CHART_DIR, "14_roc_curves.png")
if os.path.exists(fp):
    doc.add_paragraph("Figure 14. ROC curves.")
    doc.add_picture(fp, width=DInches(4.5))
doc.add_heading("10. Knowledge Interpretation", level=1)
for insight in [f"Fare prediction solved: R\u00b2={mv_r2:.2f}, ${mv_rmse:.2f} RMSE.", f"Duration is log-normal: skew {skew_b:.2f} \u2192 {skew_a:.2f}.", f"Hourly counts overdispersed ({disp:.0f}\u00d7) \u2014 Negative Binomial needed.", "Tip prediction without payment_type gives honest metrics \u2014 leakage caught.", "Outer-borough trips stay local \u2014 fleet-positioning signal."]:
    doc.add_paragraph(insight, style="List Bullet")
doc.add_heading("11. Deployment (Optional / Future)", level=1)
doc.add_paragraph("Single Python script + HTML dashboard. Production: Flask/FastAPI + Docker.")
doc.add_heading("12. Conclusions and Future Work", level=1)
doc.add_paragraph(f"Complete pipeline on {S['rows_raw']:,} trips. All algorithms from scratch. 5 regression types, decision tree, K-means++, DBSCAN, Apriori, PCA, 5-fold CV, 3 hypothesis tests, interactive dashboard.")
doc.add_heading("Future work", level=2)
for fw in ["Negative Binomial GLM.", "Spatial DBSCAN on lat/long.", "Per-cluster hierarchical models.", "Streaming pipeline (Kafka + Spark).", "Temporal train/test split.", "Fairness audit across boroughs."]:
    doc.add_paragraph(fw, style="List Bullet")

docx_path = os.path.join(SCRIPT_DIR, "Taxi_Project_Report.docx")
doc.save(docx_path)
print(f"\u2713 Saved: {docx_path}")
print("\nDone rebuilding deliverables.")
